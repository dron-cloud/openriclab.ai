from __future__ import annotations

import csv
import io
import json
import os
import time
from collections import defaultdict, deque
from pathlib import Path
from typing import Literal

import httpx
from docx import Document
from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from openpyxl import load_workbook
from pptx import Presentation
from pydantic import BaseModel, Field
from pypdf import PdfReader


# ============================================================
# Environment configuration
# ============================================================

OLLAMA_API_BASE = os.getenv(
    "OLLAMA_API_BASE",
    "https://ollama.com/api",
).rstrip("/")

OLLAMA_MODEL = os.getenv(
    "OLLAMA_MODEL",
    "gpt-oss:20b",
).strip()

OLLAMA_API_KEY = os.getenv(
    "OLLAMA_API_KEY",
    "",
).strip()

REQUEST_TIMEOUT_SECONDS = float(
    os.getenv(
        "REQUEST_TIMEOUT_SECONDS",
        "180",
    )
)

OLLAMA_NUM_PREDICT = int(
    os.getenv(
        "OLLAMA_NUM_PREDICT",
        "-1",
    )
)

OLLAMA_THINKING = os.getenv(
    "OLLAMA_THINKING",
    "low",
).strip().lower()

if OLLAMA_THINKING not in {
    "low",
    "medium",
    "high",
}:
    OLLAMA_THINKING = "low"


# ============================================================
# Public-chat limits
# ============================================================

MAX_USER_CHARACTERS = int(
    os.getenv(
        "MAX_USER_CHARACTERS",
        "6000",
    )
)

MAX_HISTORY_MESSAGES = int(
    os.getenv(
        "MAX_HISTORY_MESSAGES",
        "8",
    )
)

MAX_HISTORY_CHARACTERS_PER_MESSAGE = int(
    os.getenv(
        "MAX_HISTORY_CHARACTERS_PER_MESSAGE",
        "4000",
    )
)

RATE_LIMIT_REQUESTS = int(
    os.getenv(
        "RATE_LIMIT_REQUESTS",
        "30",
    )
)

RATE_LIMIT_WINDOW_SECONDS = int(
    os.getenv(
        "RATE_LIMIT_WINDOW_SECONDS",
        "600",
    )
)


# ============================================================
# File-upload limits
# ============================================================

MAX_FILE_BYTES = int(
    os.getenv(
        "MAX_FILE_BYTES",
        str(10 * 1024 * 1024),
    )
)

MAX_DOCUMENT_CHARACTERS = int(
    os.getenv(
        "MAX_DOCUMENT_CHARACTERS",
        "60000",
    )
)

MAX_PDF_PAGES = int(
    os.getenv(
        "MAX_PDF_PAGES",
        "80",
    )
)

MAX_SPREADSHEET_ROWS = int(
    os.getenv(
        "MAX_SPREADSHEET_ROWS",
        "3000",
    )
)

MAX_PRESENTATION_SLIDES = int(
    os.getenv(
        "MAX_PRESENTATION_SLIDES",
        "100",
    )
)


# ============================================================
# CORS
# ============================================================

ALLOWED_ORIGINS = [
    origin.strip()
    for origin in os.getenv(
        "ALLOWED_ORIGINS",
        (
            "https://dron-cloud.github.io,"
            "http://localhost:8000,"
            "http://127.0.0.1:5500,"
            "http://localhost:5500"
        ),
    ).split(",")
    if origin.strip()
]


# ============================================================
# General-purpose system prompt
# ============================================================

SYSTEM_PROMPT = """
You are AI Chat, a general-purpose AI assistant hosted on
Dr. Dara Ron's research website and powered by an external
large language model.

You can help users with:

- General knowledge
- Science and engineering
- Mathematics
- Programming and debugging
- Artificial intelligence and machine learning
- Writing and editing
- Research and technical concepts
- Education and learning
- Brainstorming
- Everyday questions and problem solving

Guidelines:

1. Answer the user's question directly.
2. Explain difficult concepts clearly.
3. Adapt the technical depth to the user's apparent expertise.
4. Use examples when useful.
5. Use equations, code, or tables when they materially improve the answer.
6. Never fabricate sources, citations, quotations, references, or factual details.
7. If you are uncertain, say what is uncertain.
8. Clearly distinguish established facts from assumptions or interpretation.
9. Do not imply that you searched the web or consulted external sources unless
   retrieval context was actually provided to you.
10. Do not claim access to real-time information unless such information is
    explicitly supplied in the conversation.
11. Keep answers concise by default but provide more detail when requested.
12. For technical questions, prioritize correctness over confidence.
13. Do not claim that you are the underlying foundation model itself.

You are a conversational assistant for general-purpose use.
""".strip()


ENGLISH_SYSTEM_PROMPT = """
Respond in clear, natural English unless the user explicitly
asks for another language.
""".strip()


# ============================================================
# Translation prompts
# ============================================================

KHMER_TO_ENGLISH_PROMPT = """
Translate the user's Khmer message into accurate, natural English.

STRICT REQUIREMENTS:
- Translate only.
- Do not answer the user's question.
- Do not add new information.
- Do not remove information.
- Do not summarize.
- Do not explain.
- Preserve the original meaning as closely as possible.
- Preserve proper names, technical terms, standards, model names,
  numbers, equations, acronyms, identifiers, URLs, code, and units.
- Preserve terms such as O-RAN, O-RU, O-DU, O-CU, 3GPP, AI, LLM,
  GPU, CPU, API, Python, HTTP, TCP/IP, 5G, and 6G.
- If a Khmer phrase is ambiguous, translate conservatively rather
  than guessing.
- Return only the English translation.
- Do not add labels such as "Translation:".
""".strip()


ENGLISH_TO_KHMER_PROMPT = """
Translate the English text into accurate, natural Khmer.

STRICT REQUIREMENTS:
- Translate only.
- Preserve the meaning of the English source exactly.
- Do not add facts.
- Do not remove facts.
- Do not independently answer or reinterpret the question.
- Do not introduce new examples or explanations.
- Use natural, modern Khmer rather than literal word-for-word translation.
- Keep important technical terms in English when that is clearer or
  technically more accurate.
- Prefer Khmer-first explanations with useful English technical terms
  in parentheses.
- Preserve proper names, standards, model names, numbers, equations,
  acronyms, identifiers, URLs, code, and units.
- Preserve established technical terms such as O-RAN, O-RU, O-DU,
  O-CU, Near-RT RIC, E2, A1, 3GPP, AI, LLM, GPU, CPU, API, Python,
  HTTP, TCP/IP, 5G, and 6G.
- Never invent a Khmer technical term if no natural equivalent exists.
- If a proper name has no standard Khmer form, keep the original name.
- Return only the Khmer translation.
- Do not add labels such as "Translation:".
""".strip()


# ============================================================
# Request models
# ============================================================

class HistoryMessage(BaseModel):
    role: Literal[
        "user",
        "assistant",
    ]

    content: str = Field(
        min_length=1,
        max_length=12000,
    )


class ChatRequest(BaseModel):
    message: str = Field(
        min_length=1,
        max_length=MAX_USER_CHARACTERS,
    )

    history: list[HistoryMessage] = Field(
        default_factory=list
    )

    language: Literal[
        "en",
        "km",
    ] = "en"


# ============================================================
# FastAPI
# ============================================================

app = FastAPI(
    title="AI Chat API",
    description=(
        "General-purpose bilingual English/Khmer AI chat "
        "interface backed by Ollama Cloud."
    ),
    version="7.1.0",
)


app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=False,
    allow_methods=[
        "GET",
        "POST",
        "OPTIONS",
    ],
    allow_headers=[
        "Content-Type",
    ],
)


# ============================================================
# Basic in-memory rate limiter
# ============================================================

request_log: dict[
    str,
    deque[float],
] = defaultdict(deque)


def get_client_ip(
    request: Request,
) -> str:
    forwarded = request.headers.get(
        "x-forwarded-for",
        "",
    )

    if forwarded:
        return forwarded.split(",")[0].strip()

    if request.client:
        return request.client.host

    return "unknown"


def enforce_rate_limit(
    request: Request,
) -> None:
    client_ip = get_client_ip(
        request
    )

    now = time.monotonic()

    timestamps = request_log[
        client_ip
    ]

    cutoff = (
        now
        - RATE_LIMIT_WINDOW_SECONDS
    )

    while (
        timestamps
        and timestamps[0] < cutoff
    ):
        timestamps.popleft()

    if (
        len(timestamps)
        >= RATE_LIMIT_REQUESTS
    ):
        raise HTTPException(
            status_code=429,
            detail=(
                "Too many requests. "
                "Please try again later."
            ),
        )

    timestamps.append(now)


# ============================================================
# Helpers
# ============================================================

def get_thinking_setting():
    if OLLAMA_MODEL.lower().startswith(
        "gpt-oss"
    ):
        return OLLAMA_THINKING

    return False


def normalize_thinking_setting(
    thinking=None,
):
    """
    GPT-OSS expects low/medium/high rather than a boolean.
    For translation calls, False is normalized to low so the
    request remains valid for GPT-OSS while minimizing reasoning.
    """
    if OLLAMA_MODEL.lower().startswith(
        "gpt-oss"
    ):
        if thinking in {
            "low",
            "medium",
            "high",
        }:
            return thinking

        return "low"

    if thinking is None:
        return False

    return bool(thinking)


def build_system_prompt() -> str:
    return (
        SYSTEM_PROMPT
        + "\n\n"
        + ENGLISH_SYSTEM_PROMPT
    )


def clean_history(
    history: list[HistoryMessage],
    current_message: str,
) -> list[dict[str, str]]:
    cleaned: list[
        dict[str, str]
    ] = []

    recent_history = history[
        -MAX_HISTORY_MESSAGES:
    ]

    for message in recent_history:
        content = (
            message.content
            .strip()[
                :MAX_HISTORY_CHARACTERS_PER_MESSAGE
            ]
        )

        if not content:
            continue

        cleaned.append(
            {
                "role":
                    message.role,

                "content":
                    content,
            }
        )

    # Remove duplicated current question if the browser
    # already included it as the final history entry.
    if (
        cleaned
        and cleaned[-1]["role"] == "user"
        and cleaned[-1]["content"].strip()
        == current_message.strip()
    ):
        cleaned.pop()

    return cleaned


# ============================================================
# Ollama client
# ============================================================

async def call_ollama(
    messages: list[dict[str, str]],
    *,
    thinking=None,
) -> str:

    thinking = normalize_thinking_setting(
        thinking
    )

    payload = {
        "model":
            OLLAMA_MODEL,

        "messages":
            messages,

        "stream":
            False,

        "think":
            thinking,
    }

    headers = {
        "Authorization":
            f"Bearer {OLLAMA_API_KEY}",

        "Content-Type":
            "application/json",

        "Accept":
            "application/json",
    }

    try:
        async with httpx.AsyncClient(
            timeout=httpx.Timeout(
                REQUEST_TIMEOUT_SECONDS
            )
        ) as client:
            response = await client.post(
                f"{OLLAMA_API_BASE}/chat",
                headers=headers,
                json=payload,
            )

    except httpx.TimeoutException as exc:
        raise HTTPException(
            status_code=504,
            detail=(
                "The AI model took too long "
                "to respond. Please try again."
            ),
        ) from exc

    except httpx.RequestError as exc:
        print(
            "Ollama connection error:",
            repr(exc),
        )

        raise HTTPException(
            status_code=502,
            detail=(
                "The AI model service is "
                "temporarily unavailable."
            ),
        ) from exc

    if response.status_code >= 400:
        error_text = response.text[:1000]

        print(
            "Ollama HTTP error:",
            response.status_code,
            error_text,
        )

        raise HTTPException(
            status_code=502,
            detail=(
                "The AI model service returned "
                "an error. Please try again."
            ),
        )

    try:
        result = response.json()

    except ValueError as exc:
        print(
            "Ollama non-JSON response:",
            response.text[:1000],
        )

        raise HTTPException(
            status_code=502,
            detail=(
                "The AI model returned an "
                "invalid response."
            ),
        ) from exc

    message = (
        result.get("message")
        if isinstance(
            result,
            dict,
        )
        else None
    )

    if not isinstance(
        message,
        dict,
    ):
        raise HTTPException(
            status_code=502,
            detail=(
                "The AI model returned an "
                "invalid message."
            ),
        )

    answer = (
        message.get(
            "content",
            "",
        )
        or ""
    ).strip()

    if not answer:
        done_reason = (
            result.get(
                "done_reason",
                "unknown",
            )
            if isinstance(
                result,
                dict,
            )
            else "unknown"
        )

        print(
            "Empty Ollama answer. "
            f"done_reason={done_reason}"
        )

        raise HTTPException(
            status_code=502,
            detail=(
                "The AI model did not return "
                "a final answer."
            ),
        )

    return answer


# ============================================================
# Ollama streaming client
# ============================================================

async def stream_ollama_events(
    messages: list[dict[str, str]],
    *,
    thinking=None,
):
    """
    Stream Ollama chat events.

    Yields:
      ("thinking", "") for hidden reasoning/progress chunks
      ("content", "<text>") for visible answer chunks

    Thinking content is intentionally never exposed to the browser.
    """

    thinking = normalize_thinking_setting(
        thinking
    )

    payload = {
        "model":
            OLLAMA_MODEL,

        "messages":
            messages,

        "stream":
            True,

        "think":
            thinking,
    }

    headers = {
        "Authorization":
            f"Bearer {OLLAMA_API_KEY}",

        "Content-Type":
            "application/json",

        "Accept":
            "application/x-ndjson",
    }

    try:
        async with httpx.AsyncClient(
            timeout=httpx.Timeout(
                REQUEST_TIMEOUT_SECONDS
            )
        ) as client:

            async with client.stream(
                "POST",
                f"{OLLAMA_API_BASE}/chat",
                headers=headers,
                json=payload,
            ) as response:

                if response.status_code >= 400:
                    error_bytes = (
                        await response.aread()
                    )

                    error_text = (
                        error_bytes.decode(
                            "utf-8",
                            errors="replace",
                        )[:1000]
                    )

                    print(
                        "Ollama streaming HTTP error:",
                        response.status_code,
                        error_text,
                    )

                    raise RuntimeError(
                        "The AI model service returned "
                        "an error. Please try again."
                    )

                async for line in response.aiter_lines():
                    if not line:
                        continue

                    try:
                        event = json.loads(
                            line
                        )
                    except json.JSONDecodeError:
                        print(
                            "Skipping invalid Ollama "
                            "stream line:",
                            line[:500],
                        )
                        continue

                    if not isinstance(
                        event,
                        dict,
                    ):
                        continue

                    message = event.get(
                        "message"
                    )

                    if isinstance(
                        message,
                        dict,
                    ):
                        thinking_chunk = (
                            message.get(
                                "thinking",
                                "",
                            )
                            or ""
                        )

                        if thinking_chunk:
                            # Signal progress without exposing
                            # private reasoning content.
                            yield (
                                "thinking",
                                "",
                            )

                        content = (
                            message.get(
                                "content",
                                "",
                            )
                            or ""
                        )

                        if content:
                            yield (
                                "content",
                                content,
                            )

                    if event.get(
                        "done"
                    ) is True:
                        break

    except httpx.TimeoutException as exc:
        print(
            "Ollama streaming timeout:",
            repr(exc),
        )

        raise RuntimeError(
            "The AI model took too long "
            "to respond. Please try again."
        ) from exc

    except httpx.RequestError as exc:
        print(
            "Ollama streaming connection error:",
            repr(exc),
        )

        raise RuntimeError(
            "The AI model service is "
            "temporarily unavailable."
        ) from exc


def ndjson_event(
    event_type: str,
    **data,
) -> str:
    return (
        json.dumps(
            {
                "type":
                    event_type,

                **data,
            },
            ensure_ascii=False,
        )
        + "\n"
    )


# ============================================================
# Translation functions
# ============================================================

async def translate_khmer_to_english(
    khmer_text: str,
) -> str:

    messages = [
        {
            "role":
                "system",

            "content":
                KHMER_TO_ENGLISH_PROMPT,
        },
        {
            "role":
                "user",

            "content":
                khmer_text,
        },
    ]

    return await call_ollama(
        messages,
        thinking=False,
    )


async def translate_english_to_khmer(
    english_text: str,
) -> str:

    messages = [
        {
            "role":
                "system",

            "content":
                ENGLISH_TO_KHMER_PROMPT,
        },
        {
            "role":
                "user",

            "content":
                english_text,
        },
    ]

    return await call_ollama(
        messages,
        thinking=False,
    )



# ============================================================
# File extraction helpers
# ============================================================

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".log",
    ".csv", ".tsv", ".json", ".jsonl", ".ndjson",
    ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".html", ".htm", ".css", ".scss", ".sass", ".less",
    ".js", ".mjs", ".cjs", ".jsx", ".ts", ".tsx",
    ".py", ".ipynb", ".java", ".c", ".h", ".cpp", ".hpp",
    ".cc", ".cxx", ".cs", ".go", ".rs", ".rb", ".php",
    ".swift", ".kt", ".kts", ".scala", ".sh", ".bash",
    ".zsh", ".fish", ".ps1", ".sql", ".r", ".m",
    ".tex", ".bib", ".gradle", ".properties", ".env",
    ".gitignore", ".dockerfile",
}

BLOCKED_EXTENSIONS = {
    ".exe", ".dll", ".so", ".dylib", ".bin", ".app",
    ".msi", ".dmg", ".pkg", ".apk", ".ipa",
    ".jar", ".war",
    ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz",
    ".iso",
}

IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff",
    ".webp", ".heic", ".heif", ".svg",
}

LEGACY_OFFICE_EXTENSIONS = {
    ".doc", ".xls", ".ppt",
}


def truncate_document_text(
    text: str,
) -> str:
    text = (
        text
        .replace("\x00", "")
        .strip()
    )

    if (
        len(text)
        <= MAX_DOCUMENT_CHARACTERS
    ):
        return text

    return (
        text[
            :MAX_DOCUMENT_CHARACTERS
        ]
        + "\n\n[Document truncated because it exceeded the "
        "configured text limit.]"
    )


def decode_text_bytes(
    data: bytes,
) -> str:
    """
    Decode common text encodings without adding a heavy charset detector.
    """
    for encoding in (
        "utf-8-sig",
        "utf-8",
        "utf-16",
        "utf-16-le",
        "utf-16-be",
        "latin-1",
    ):
        try:
            return data.decode(
                encoding
            )
        except UnicodeDecodeError:
            continue

    return data.decode(
        "utf-8",
        errors="replace",
    )


def extract_pdf_text(
    data: bytes,
) -> str:
    try:
        reader = PdfReader(
            io.BytesIO(data)
        )
    except Exception as exc:
        raise HTTPException(
            status_code=415,
            detail=(
                "The PDF could not be opened or is invalid."
            ),
        ) from exc

    if reader.is_encrypted:
        try:
            unlocked = reader.decrypt("")
        except Exception:
            unlocked = 0

        if not unlocked:
            raise HTTPException(
                status_code=415,
                detail=(
                    "Password-protected PDFs are not supported."
                ),
            )

    parts: list[str] = []

    for page_number, page in enumerate(
        reader.pages[:MAX_PDF_PAGES],
        start=1,
    ):
        try:
            page_text = (
                page.extract_text()
                or ""
            ).strip()
        except Exception as exc:
            print(
                f"PDF page {page_number} extraction error:",
                repr(exc),
            )
            page_text = ""

        if page_text:
            parts.append(
                f"\n--- PDF PAGE {page_number} ---\n"
                f"{page_text}"
            )

        if (
            sum(
                len(part)
                for part in parts
            )
            >= MAX_DOCUMENT_CHARACTERS
        ):
            break

    text = "\n".join(
        parts
    ).strip()

    if not text:
        raise HTTPException(
            status_code=422,
            detail=(
                "No extractable text was found in this PDF. "
                "It may be a scanned/image-only PDF. OCR is not "
                "enabled on this public endpoint."
            ),
        )

    return truncate_document_text(
        text
    )


def extract_docx_text(
    data: bytes,
) -> str:
    try:
        document = Document(
            io.BytesIO(data)
        )
    except Exception as exc:
        raise HTTPException(
            status_code=415,
            detail=(
                "The DOCX file could not be opened."
            ),
        ) from exc

    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = (
            paragraph.text
            or ""
        ).strip()

        if text:
            parts.append(text)

    for table_index, table in enumerate(
        document.tables,
        start=1,
    ):
        parts.append(
            f"\n--- TABLE {table_index} ---"
        )

        for row in table.rows:
            values = [
                (
                    cell.text
                    or ""
                ).strip()
                for cell in row.cells
            ]

            if any(values):
                parts.append(
                    " | ".join(values)
                )

    text = "\n".join(
        parts
    ).strip()

    if not text:
        raise HTTPException(
            status_code=422,
            detail=(
                "No extractable text was found in the DOCX file."
            ),
        )

    return truncate_document_text(
        text
    )


def extract_xlsx_text(
    data: bytes,
) -> str:
    try:
        workbook = load_workbook(
            filename=io.BytesIO(data),
            read_only=True,
            data_only=True,
        )
    except Exception as exc:
        raise HTTPException(
            status_code=415,
            detail=(
                "The Excel workbook could not be opened."
            ),
        ) from exc

    parts: list[str] = []
    row_count = 0

    try:
        for worksheet in workbook.worksheets:
            parts.append(
                f"\n--- SHEET: {worksheet.title} ---"
            )

            for row in worksheet.iter_rows(
                values_only=True
            ):
                values = [
                    ""
                    if value is None
                    else str(value)
                    for value in row
                ]

                if any(
                    value.strip()
                    for value in values
                ):
                    parts.append(
                        "\t".join(values)
                    )

                    row_count += 1

                if (
                    row_count
                    >= MAX_SPREADSHEET_ROWS
                ):
                    parts.append(
                        "\n[Spreadsheet row limit reached.]"
                    )
                    break

            if (
                row_count
                >= MAX_SPREADSHEET_ROWS
            ):
                break
    finally:
        workbook.close()

    text = "\n".join(
        parts
    ).strip()

    if not text:
        raise HTTPException(
            status_code=422,
            detail=(
                "No readable cell content was found in the workbook."
            ),
        )

    return truncate_document_text(
        text
    )


def extract_pptx_text(
    data: bytes,
) -> str:
    try:
        presentation = Presentation(
            io.BytesIO(data)
        )
    except Exception as exc:
        raise HTTPException(
            status_code=415,
            detail=(
                "The PowerPoint file could not be opened."
            ),
        ) from exc

    parts: list[str] = []

    slides = presentation.slides[
        :MAX_PRESENTATION_SLIDES
    ]

    for slide_number, slide in enumerate(
        slides,
        start=1,
    ):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            if hasattr(
                shape,
                "text"
            ):
                text = (
                    shape.text
                    or ""
                ).strip()

                if text:
                    slide_parts.append(
                        text
                    )

            if getattr(
                shape,
                "has_table",
                False,
            ):
                table = shape.table

                for row in table.rows:
                    values = [
                        (
                            cell.text
                            or ""
                        ).strip()
                        for cell in row.cells
                    ]

                    if any(values):
                        slide_parts.append(
                            " | ".join(values)
                        )

        if slide_parts:
            parts.append(
                f"\n--- SLIDE {slide_number} ---\n"
                + "\n".join(
                    slide_parts
                )
            )

    text = "\n".join(
        parts
    ).strip()

    if not text:
        raise HTTPException(
            status_code=422,
            detail=(
                "No extractable text was found in the presentation."
            ),
        )

    return truncate_document_text(
        text
    )


def extract_csv_or_tsv_text(
    data: bytes,
    extension: str,
) -> str:
    text = decode_text_bytes(
        data
    )

    delimiter = (
        "\t"
        if extension == ".tsv"
        else ","
    )

    reader = csv.reader(
        io.StringIO(text),
        delimiter=delimiter,
    )

    rows: list[str] = []

    for index, row in enumerate(
        reader,
        start=1,
    ):
        if (
            index
            > MAX_SPREADSHEET_ROWS
        ):
            rows.append(
                "[CSV/TSV row limit reached.]"
            )
            break

        rows.append(
            " | ".join(
                str(value)
                for value in row
            )
        )

    return truncate_document_text(
        "\n".join(rows)
    )


def extract_uploaded_file_text(
    filename: str,
    content_type: str,
    data: bytes,
) -> str:
    """
    Extract text from the broad set of document formats supported
    by this public chat endpoint.

    This intentionally does not execute, unpack, or render uploaded
    user files.
    """

    extension = (
        Path(filename)
        .suffix
        .lower()
    )

    lower_name = filename.lower()

    if extension in BLOCKED_EXTENSIONS:
        raise HTTPException(
            status_code=415,
            detail=(
                "Executable and archive files are not supported."
            ),
        )

    if extension in IMAGE_EXTENSIONS:
        raise HTTPException(
            status_code=415,
            detail=(
                "Image uploads require OCR/vision processing, which "
                "is not enabled on this endpoint yet."
            ),
        )

    if extension in LEGACY_OFFICE_EXTENSIONS:
        raise HTTPException(
            status_code=415,
            detail=(
                "Legacy .doc, .xls, and .ppt files are not supported. "
                "Please save them as DOCX, XLSX, or PPTX."
            ),
        )

    if extension == ".pdf":
        return extract_pdf_text(
            data
        )

    if extension == ".docx":
        return extract_docx_text(
            data
        )

    if extension in {
        ".xlsx",
        ".xlsm",
        ".xltx",
        ".xltm",
    }:
        return extract_xlsx_text(
            data
        )

    if extension in {
        ".pptx",
        ".ppsx",
        ".potx",
    }:
        return extract_pptx_text(
            data
        )

    if extension in {
        ".csv",
        ".tsv",
    }:
        return extract_csv_or_tsv_text(
            data,
            extension,
        )

    if (
        extension in TEXT_EXTENSIONS
        or lower_name in {
            "dockerfile",
            "makefile",
            "requirements.txt",
            "readme",
            "license",
        }
        or (
            content_type
            and (
                content_type.startswith(
                    "text/"
                )
                or content_type in {
                    "application/json",
                    "application/xml",
                    "application/javascript",
                    "application/sql",
                    "application/x-yaml",
                    "application/yaml",
                }
            )
        )
    ):
        return truncate_document_text(
            decode_text_bytes(
                data
            )
        )

    # Last-resort safe text attempt for unknown, non-binary-looking files.
    # Reject if the sample contains too many NUL bytes.
    sample = data[:4096]

    if (
        sample
        and sample.count(
            b"\x00"
        )
        <= max(
            1,
            len(sample) // 100,
        )
    ):
        guessed_text = truncate_document_text(
            decode_text_bytes(
                data
            )
        )

        if guessed_text.strip():
            return guessed_text

    raise HTTPException(
        status_code=415,
        detail=(
            "Unsupported file type. Supported formats include PDF, "
            "DOCX, XLSX/XLSM, PPTX, CSV/TSV, Markdown, JSON, XML, "
            "YAML, HTML, source code, logs, configuration files, "
            "and ordinary text files."
        ),
    )


def build_document_prompt(
    *,
    filename: str,
    document_text: str,
    question: str,
) -> str:
    return f"""
The user uploaded a document.

FILE:
{filename}

DOCUMENT CONTENT:
{document_text}

USER QUESTION:
{question}

Instructions:
- Answer the user's question using the uploaded document as the primary source.
- Do not claim the document contains information that is not present.
- If the question is document-specific and the document does not provide enough
  information, clearly say that the document does not contain enough information.
- Preserve technical terminology, equations, identifiers, values, and names
  exactly when relevant.
""".strip()


# ============================================================
# Routes
# ============================================================

@app.get("/")
async def root() -> dict:
    return {
        "service":
            "AI Chat",

        "status":
            "ok",

        "model":
            OLLAMA_MODEL,

        "languages": [
            "en",
            "km",
        ],

        "streaming":
            True,

        "stream_format":
            "ndjson",

        "khmer_pipeline":
            "translation-assisted",
    }


@app.get("/api/health")
async def health() -> dict:
    return {
        "status": (
            "ok"
            if OLLAMA_API_KEY
            else "missing_api_key"
        ),

        "model":
            OLLAMA_MODEL,

        "ollama_api_base":
            OLLAMA_API_BASE,

        "num_predict":
            OLLAMA_NUM_PREDICT,

        "thinking":
            get_thinking_setting(),

        "languages": [
            "en",
            "km",
        ],

        "streaming":
            True,

        "stream_format":
            "ndjson",

        "khmer_pipeline":
            "km->en->AI->km",
    }


@app.post("/api/chat")
async def chat(
    request_body: ChatRequest,
    request: Request,
):
    # normal text-only streaming chat
    ...


@app.post("/api/chat/file")
async def chat_with_file(
    request: Request,
    message: str = Form(""),
    language: str = Form("en"),
    file: UploadFile = File(...),
):
    # document-upload streaming chat
    ...

    # --------------------------------------------------------
    # Validate request before streaming begins
    # --------------------------------------------------------

    enforce_rate_limit(
        request
    )

    if not OLLAMA_API_KEY:
        raise HTTPException(
            status_code=500,
            detail=(
                "The AI service is not "
                "configured correctly."
            ),
        )

    user_message = (
        message
        .strip()
    )

    if not user_message:
        raise HTTPException(
            status_code=422,
            detail=(
                "Please enter a question or instruction "
                "for the uploaded file."
            ),
        )

    if (
        len(user_message)
        > MAX_USER_CHARACTERS
    ):
        raise HTTPException(
            status_code=413,
            detail=(
                "Your message is too long. "
                f"Please keep it below "
                f"{MAX_USER_CHARACTERS} characters."
            ),
        )

    if language not in {
        "en",
        "km",
    }:
        raise HTTPException(
            status_code=422,
            detail="Unsupported language.",
        )

    filename = (
        file.filename
        or "uploaded-file"
    )

    content_type = (
        file.content_type
        or "application/octet-stream"
    )

    data = await file.read()

    try:
        await file.close()
    except Exception:
        pass

    if not data:
        raise HTTPException(
            status_code=422,
            detail=(
                "The uploaded file is empty."
            ),
        )

    if (
        len(data)
        > MAX_FILE_BYTES
    ):
        max_mb = (
            MAX_FILE_BYTES
            / (1024 * 1024)
        )

        raise HTTPException(
            status_code=413,
            detail=(
                "The uploaded file is too large. "
                f"Maximum size is {max_mb:.0f} MB."
            ),
        )

    print(
        "AI Chat file upload:",
        {
            "filename":
                filename,
            "content_type":
                content_type,
            "size_bytes":
                len(data),
            "language":
                language,
        },
    )

    document_text = (
        extract_uploaded_file_text(
            filename,
            content_type,
            data,
        )
    )

    if not document_text.strip():
        raise HTTPException(
            status_code=422,
            detail=(
                "No readable text could be extracted "
                "from the uploaded file."
            ),
        )

    # --------------------------------------------------------
    # Browser-facing NDJSON generator
    # --------------------------------------------------------

    async def browser_stream():
        try:
            yield ndjson_event(
                "start",
                language=language,
                filename=filename,
            )

            # ====================================================
            # KHMER MODE
            # ====================================================

            if language == "km":

                yield ndjson_event(
                    "status",
                    stage="translating_question",
                )

                english_question = (
                    await translate_khmer_to_english(
                        user_message
                    )
                )

                yield ndjson_event(
                    "status",
                    stage="reasoning",
                )

                english_messages = [
                    {
                        "role":
                            "system",

                        "content":
                            build_system_prompt(),
                    },
                    {
                        "role":
                            "user",

                        "content":
                            build_document_prompt(
                                filename=filename,
                                document_text=document_text,
                                question=english_question,
                            ),
                    },
                ]

                english_answer = (
                    await call_ollama(
                        english_messages,
                        thinking=get_thinking_setting(),
                    )
                )

                yield ndjson_event(
                    "status",
                    stage="translating_answer",
                )

                translation_messages = [
                    {
                        "role":
                            "system",

                        "content":
                            ENGLISH_TO_KHMER_PROMPT,
                    },
                    {
                        "role":
                            "user",

                        "content":
                            english_answer,
                    },
                ]

                content_seen = False

                async for (
                    event_type,
                    content,
                ) in stream_ollama_events(
                    translation_messages,
                    thinking="low",
                ):
                    if (
                        event_type
                        == "thinking"
                    ):
                        yield ndjson_event(
                            "status",
                            stage="translating_answer",
                        )
                        continue

                    if (
                        event_type
                        == "content"
                        and content
                    ):
                        content_seen = True

                        yield ndjson_event(
                            "chunk",
                            content=content,
                        )

                if not content_seen:
                    raise RuntimeError(
                        "The AI model returned "
                        "an empty Khmer response."
                    )

                yield ndjson_event(
                    "done",
                    language="km",
                    model=OLLAMA_MODEL,
                    pipeline="file->km->en->AI->km",
                    filename=filename,
                )

                return

            # ====================================================
            # ENGLISH MODE
            # ====================================================

            messages = [
                {
                    "role":
                        "system",

                    "content":
                        build_system_prompt(),
                },
                {
                    "role":
                        "user",

                    "content":
                        build_document_prompt(
                            filename=filename,
                            document_text=document_text,
                            question=user_message,
                        ),
                },
            ]

            content_seen = False

            yield ndjson_event(
                "status",
                stage="reasoning",
            )

            async for (
                event_type,
                content,
            ) in stream_ollama_events(
                messages,
                thinking=get_thinking_setting(),
            ):
                if (
                    event_type
                    == "thinking"
                ):
                    yield ndjson_event(
                        "status",
                        stage="reasoning",
                    )
                    continue

                if (
                    event_type
                    == "content"
                    and content
                ):
                    content_seen = True

                    yield ndjson_event(
                        "chunk",
                        content=content,
                    )

            if not content_seen:
                raise RuntimeError(
                    "The AI model returned "
                    "an empty response."
                )

            yield ndjson_event(
                "done",
                language="en",
                model=OLLAMA_MODEL,
                pipeline="file->AI",
                filename=filename,
            )

        except Exception as exc:
            print(
                "AI Chat file-streaming error:",
                repr(exc),
            )

            public_message = (
                str(exc)
                .strip()
                or (
                    "The AI service encountered an "
                    "error while processing the file."
                )
            )

            yield ndjson_event(
                "error",
                message=public_message,
            )

    return StreamingResponse(
        browser_stream(),
        media_type=(
            "application/x-ndjson; "
            "charset=utf-8"
        ),
        headers={
            "Cache-Control":
                "no-cache, no-transform",

            "X-Accel-Buffering":
                "no",

            "Content-Type":
                "application/x-ndjson; charset=utf-8",
        },
    )
